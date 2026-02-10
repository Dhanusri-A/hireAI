import os
import json
import requests
from fastapi import HTTPException, UploadFile

# For PDF parsing
from io import BytesIO
from pypdf import PdfReader
from docx import Document
from pptx import Presentation
from pptx.shapes.base import BaseShape
from pptx.shapes.autoshape import Shape







class ResumeParser:
    def __init__(self):
        self.api_key = os.getenv("GROQ_API_KEY")

        if not self.api_key:
            try:
                with open("groq_key.txt", "r") as f:
                    self.api_key = f.read().strip()
            except FileNotFoundError:
                self.api_key = None

        self.api_url = "https://api.groq.com/openai/v1/chat/completions"


    async def resume_text_extract(self,file: UploadFile) -> str:
        """
        Extract text from an uploaded resume file (PDF, DOCX, PPTX).
        Accepts FastAPI UploadFile.
        Returns plain text.
        """

        filename = (file.filename or "").lower()
        content:bytes = await file.read()
        file_like = BytesIO(content)

        # ---------------- PDF ----------------
        if filename.endswith(".pdf"):
            reader = PdfReader(file_like)
            text_parts = []

            for page in reader.pages:
                page_text = page.extract_text()
                if page_text:
                    text_parts.append(page_text)

            return "\n".join(text_parts).strip()

        # ---------------- DOCX ----------------
        elif filename.endswith(".docx"):
            doc = Document(file_like)
            text_parts = [p.text for p in doc.paragraphs if p.text.strip()]
            return "\n".join(text_parts).strip()

        # ---------------- PPTX ----------------
        elif filename.endswith(".pptx"):       
            prs = Presentation(file_like)
            text_parts: list[str] = []

            for slide in prs.slides:
                for shape in slide.shapes:
                    if isinstance(shape, Shape) and shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            if para.text.strip():
                                text_parts.append(para.text)

            return "\n".join(text_parts).strip()

        # ---------------- Unsupported ----------------
        else:
            raise ValueError(
                "Unsupported file type. Only PDF, DOCX, and PPTX are allowed."
            )

    async def parse_resume(self, file: UploadFile) -> dict:

        resume_text: str = await self.resume_text_extract(file)

        if not self.api_key:
            raise HTTPException(
                status_code=500,
                detail="Groq API key not configured",
            )

        prompt = self._build_prompt(resume_text)

        payload = {
            "model": "llama-3.1-8b-instant",
            "messages": [
                {
                    "role": "system",
                    "content": (
                        "You are a resume parsing expert. "
                        "Extract structured information from resumes. "
                        "Return ONLY valid JSON. No markdown. No explanations."
                    ),
                },
                {"role": "user", "content": prompt},
            ],
            "temperature": 0.1,
            "max_tokens": 2000,
        }

        headers = {
            "Authorization": f"Bearer {self.api_key}",
            "Content-Type": "application/json",
        }

        try:
            response = requests.post(
                self.api_url,
                headers=headers,
                json=payload,
                timeout=30,
            )

            if response.status_code != 200:
                raise HTTPException(
                    status_code=response.status_code,
                    detail=response.text,
                )

            content = response.json()["choices"][0]["message"]["content"]

            try:
                return json.loads(content)
            except json.JSONDecodeError:
                raise HTTPException(
                    status_code=500,
                    detail="LLM returned invalid JSON",
                )

        except requests.RequestException as exc:
            raise HTTPException(
                status_code=500,
                detail=f"Groq API request failed: {exc}",
            )

    def _build_prompt(self, text: str) -> str:
        return """

You MUST respond with ONLY valid JSON.
NO markdown. NO comments. NO explanations.

Return JSON in EXACTLY this format:

{
  "first_name": "",
  "last_name": "",
  "email": "",
  "title": "",
  "phone": "",
  "location": "",
  "skills": "",
  "profile_summary": "",
  "total_years_experience": "",
  "notice_period": "",
  "expected_salary": "",
  "preferred_mode": "",
  "profiles": {
    "linkedin": "",
    "github": ""
  },
  "languages": {
    "english": "",
    "hindi": ""
  },
  "education": {
    "institution_name": "",
    "degree": "",
    "field_of_study": "",
    "start_year": "",
    "end_year": "",
    "gpa": "",
    "honors": ""
  },
  "work_experience": {
    "company_name": "",
    "job_title": "",
    "start_date": "",
    "end_date": "",
    "location": "",
    "description": ""
  }
}

STRICT RULES:
- Missing values must be empty strings or empty arrays
- Do NOT invent data
- Do NOT use the word "Unknown"
- Only extract what exists

Resume Text:
""" + text



resume_parser = ResumeParser()
