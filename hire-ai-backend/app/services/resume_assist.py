import os
import uuid
import tempfile
import re
import zipfile

import fitz  # PyMuPDF
from pypdf import PdfReader
from docx import Document
from pptx import Presentation

from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
from reportlab.lib.styles import getSampleStyleSheet


UPLOAD_FOLDER = tempfile.gettempdir()

def convert_docx_to_pdf(docx_path: str, pdf_path: str) -> bool:
    try:
        doc = Document(docx_path)
        styles = getSampleStyleSheet()
        story = []

        for p in doc.paragraphs:
            if p.text.strip():
                story.append(Paragraph(p.text, styles["Normal"]))
                story.append(Spacer(1, 12))

        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    if cell.text.strip():
                        story.append(Paragraph(cell.text, styles["Normal"]))
                        story.append(Spacer(1, 12))

        pdf = SimpleDocTemplate(pdf_path, pagesize=letter)
        pdf.build(story)
        return True

    except Exception as e:
        print(f"DOCX → PDF failed: {e}")
        return False


def convert_pptx_to_pdf(pptx_path: str, pdf_path: str) -> bool:
    try:
        prs = Presentation(pptx_path)
        pdf = fitz.open()

        for slide in prs.slides:
            page = pdf.new_page()
            rect = fitz.Rect(50, 50, 550, 750)

            slide_text = ""
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text += shape.text + "\n"

            page.insert_textbox(rect, slide_text, fontsize=11)

        pdf.save(pdf_path)
        pdf.close()
        return True

    except Exception as e:
        print(f"PPTX → PDF failed: {e}")
        return False



def extract_text_from_pdf(pdf_path: str) -> str:
    doc = fitz.open(pdf_path)
    text_parts: list[str] = []

    for page in doc:
        page_text = page.get_text()
        if isinstance(page_text, str):
            text_parts.append(page_text)

    doc.close()
    return "".join(text_parts)


def extract_text_from_legacy_doc(doc_path: str) -> str:
    """
    Best-effort extractor for legacy binary .doc files.
    It pulls readable text runs and filters noisy binary fragments.
    """
    with open(doc_path, "rb") as file_pointer:
        raw = file_pointer.read()

    decoded_text_candidates = []

    for encoding in ("utf-16le", "utf-8", "latin-1"):
        try:
            decoded_text_candidates.append(raw.decode(encoding, errors="ignore"))
        except Exception:
            continue

    combined = "\n".join(decoded_text_candidates)
    cleaned = re.sub(r"[^\x09\x0A\x0D\x20-\x7E]", " ", combined)
    cleaned = re.sub(r"\s+", " ", cleaned)

    fragments = re.findall(r"[A-Za-z0-9@._%+\-/(),:&' ]{4,}", cleaned)
    lines = [fragment.strip() for fragment in fragments if len(fragment.strip()) >= 4]

    unique_lines = []
    seen = set()
    for line in lines:
        key = line.lower()
        if key not in seen:
            seen.add(key)
            unique_lines.append(line)

    return "\n".join(unique_lines)


def extract_text_from_doc(doc_path: str) -> str:
    """
    Supports both:
    - .docx content uploaded with .doc extension
    - legacy binary .doc files (best-effort text extraction)
    """
    if zipfile.is_zipfile(doc_path):
        try:
            doc = Document(doc_path)
            parts = [paragraph.text for paragraph in doc.paragraphs if paragraph.text.strip()]

            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        if cell.text.strip():
                            parts.append(cell.text.strip())

            text = "\n".join(parts).strip()
            if text:
                return text
        except Exception:
            pass

    return extract_text_from_legacy_doc(doc_path)



def process_uploaded_file(file) -> str:
    """
    Accepts uploaded file (PDF / DOC / DOCX / PPTX)
    Returns extracted text only
    """

    original_ext = os.path.splitext(file.filename)[1].lower()
    temp_id = str(uuid.uuid4())

    original_path = os.path.join(
        UPLOAD_FOLDER, f"input_{temp_id}{original_ext}"
    )

    with open(original_path, "wb") as f:
        f.write(file.file.read())

    # PDF → extract directly
    if original_ext == ".pdf":
        return extract_text_from_pdf(original_path)

    if original_ext == ".doc":
        extracted = extract_text_from_doc(original_path).strip()
        if not extracted:
            raise RuntimeError("Could not extract text from DOC file")
        return extracted

    # DOCX / PPTX → convert → extract
    pdf_path = os.path.join(
        UPLOAD_FOLDER, f"converted_{temp_id}.pdf"
    )

    if original_ext == ".docx":
        success = convert_docx_to_pdf(original_path, pdf_path)

    elif original_ext == ".pptx":
        success = convert_pptx_to_pdf(original_path, pdf_path)

    else:
        raise ValueError("Unsupported file type")

    if not success:
        raise RuntimeError("Conversion to PDF failed")

    return extract_text_from_pdf(pdf_path)



