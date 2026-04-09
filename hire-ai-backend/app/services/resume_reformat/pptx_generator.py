import os
import tempfile
import traceback

from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.dml.color import RGBColor

SYNPULSE_ORANGE = RGBColor(0xF3, 0x73, 0x20)
SYNPULSE_MAROON = RGBColor(0x73, 0x23, 0x3C)
TEXT_DARK       = RGBColor(0x33, 0x33, 0x33)

SHAPE_MAP = {
    "Candidate_Name":      "TextBox 1",
    "Candidate_Title":     "Candidate_Title",
    "About_Summary":       "About_Summary",
    "Education_Content":   "Education_Content",
    "Core_Skills_Content": "Core_Skills_Content",
    "Languages_Content":   "Languages_Content",
    "Hobbies_Content":     "Hobbies_Content",
    "Experience_Content":  "TextBox 5",
    "Quote":               "Quote",
}



def set_core_skills_shape(slide, skills: dict):
    shape_name = SHAPE_MAP["Core_Skills_Content"]

    for shape in slide.shapes:
        if shape.name != shape_name:
            continue
        if not shape.has_text_frame:
            return

        tf = shape.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.NONE

        first_para = True

        for cat in ['project_management', 'technical', 'devops', 'domains']:
            items = skills.get(cat, [])
            if not items:
                continue

            title = cat.replace('_', ' ').title()

            para = tf.paragraphs[0] if first_para else tf.add_paragraph()
            first_para = False

            # 🔴 Title (Maroon + Bold)
            run = para.add_run()
            run.text = f"{title}: "
            run.font.size = Pt(8)
            run.font.bold = True
            run.font.color.rgb = SYNPULSE_MAROON

            # ⚫ Items (Dark)
            run = para.add_run()
            run.text = ", ".join(map(str, items))
            run.font.size = Pt(8)
            run.font.bold = False
            run.font.color.rgb = TEXT_DARK

        return

    print(f"[ERROR] Shape not found: '{shape_name}'")

def format_education(edu_list: list) -> str:
    lines = []
    for e in edu_list[:4]:
        parts = []
        degree = e.get("degree","").strip()
        inst   = e.get("institution","").strip()
        year   = e.get("year","").strip()
        spec   = e.get("specialization","").strip()
        if degree: parts.append(degree)
        if inst:   parts.append(inst)
        if year:   parts.append(f"({year})")
        if parts:
            line = " – ".join(parts)
            if spec: line += f" | {spec}"
            lines.append(line)
    return "\n".join(lines) if lines else "—"


def format_experience_brief(exp_list: list) -> str:
    """Role | Company header + first 3 responsibilities max per job."""
    lines = []
    for exp in exp_list:
        role     = exp.get("role","").strip()
        company  = exp.get("company","").strip()
        duration = exp.get("duration","").strip()
        header   = role
        if company:  header += f" | {company}"
        if duration: header += f" ({duration})"
        if header: lines.append(header)
        for resp in exp.get("responsibilities", [])[:3]:
            cleaned = str(resp).strip()
            if cleaned: lines.append(f"• {cleaned}")
        lines.append("")
    while lines and not lines[-1]: lines.pop()
    return "\n".join(lines) if lines else "—"


# ─── Set plain text into a shape ──────────────────────────────────────────────

def set_shape_text(slide, key, text, size=9, bold=False, color=TEXT_DARK):
    shape_name = SHAPE_MAP.get(key)
    if not shape_name:
        print(f"[WARN] No shape map for '{key}'"); return
    for shape in slide.shapes:
        if shape.name == shape_name:
            if not shape.has_text_frame:
                print(f"[WARN] No text frame: '{shape_name}'"); return
            tf = shape.text_frame
            tf.clear(); tf.word_wrap = True; tf.auto_size = MSO_AUTO_SIZE.NONE
            for i, line in enumerate(str(text).split("\n")):
                para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                run = para.add_run()
                run.text = line; run.font.size = Pt(size)
                run.font.bold = bold; run.font.color.rgb = color
            print(f"[OK] {key} → '{str(text)[:70]}'"); return
    print(f"[ERROR] Shape not found: '{shape_name}'")


def set_experience_shape(slide, exp_list):
    """
    Set experience TextBox 5 with bold maroon headers per job
    and dark-coloured bullet points underneath.
    Max 3 bullets per job to fit layout.
    """
    shape_name = SHAPE_MAP["Experience_Content"]
    for shape in slide.shapes:
        if shape.name != shape_name: continue
        if not shape.has_text_frame: return
        tf = shape.text_frame
        tf.clear(); tf.word_wrap = True; tf.auto_size = MSO_AUTO_SIZE.NONE

        first_para = True
        for exp in exp_list:
            role     = exp.get("role","").strip()
            company  = exp.get("company","").strip()
            duration = exp.get("duration","").strip()
            header   = role
            if company:  header += f" | {company}"
            if duration: header += f" ({duration})"

            # Bold maroon header
            para = tf.paragraphs[0] if first_para else tf.add_paragraph()
            first_para = False
            run = para.add_run()
            run.text = header; run.font.size = Pt(9)
            run.font.bold = True; run.font.color.rgb = SYNPULSE_MAROON

            # Up to 3 bullet lines in dark colour
            for resp in exp.get("responsibilities", [])[:5]:
                cleaned = str(resp).strip()
                if not cleaned: continue
                para = tf.add_paragraph()
                run  = para.add_run()
                run.text = f"• {cleaned}"; run.font.size = Pt(8)
                run.font.bold = False; run.font.color.rgb = TEXT_DARK

            # Blank spacer between jobs
            tf.add_paragraph()
        return
    print(f"[ERROR] Shape not found: '{shape_name}'")


# ─── Move a shape's top position ──────────────────────────────────────────────

def _move_shape_top(slide, shape_name, new_top_inches):
    for s in slide.shapes:
        if s.name == shape_name:
            s.top = int(new_top_inches * 914400)
            return


# ─── Main generator ───────────────────────────────────────────────────────────

def generate_synpulse_pptx(data: dict) -> str:
    if not data: raise ValueError("No data for PPTX")
    BASE_DIR      = os.path.dirname(os.path.abspath(__file__))
    template_path = os.path.join(BASE_DIR, "format-pptx.pptx")
    if not os.path.isfile(template_path):
        raise FileNotFoundError(f"Template not found: {template_path}")

    print("\n=== STARTING PPTX GENERATION ===")
    try:
        prs  = Presentation(template_path)
        slide = prs.slides[0]
        personal = data.get("personalInfo", {})
        print(personal)

        # ── Move name & title UP so they don't overlap the connector line ─────
        # Original: TextBox 1 (name) top=2.273", Candidate_Title top=2.743"
        # Connector line is at top=2.977"  — shift both up by ~0.45"
        _move_shape_top(slide, "TextBox 1",        1.82)
        _move_shape_top(slide, "Candidate_Title",   2.29)

        # ── Name ──────────────────────────────────────────────────────────────
        name = personal.get("name","").strip()
        set_shape_text(slide, "Candidate_Name",
                       name.upper() if name else "—",
                       size=14, bold=True, color=SYNPULSE_MAROON)

        # ── Title ─────────────────────────────────────────────────────────────
        title = personal.get("title","").strip()
        set_shape_text(slide, "Candidate_Title",
                       title if title else "—",
                       size=10, color=SYNPULSE_ORANGE)

        # ── Summary ───────────────────────────────────────────────────────────
        summary = personal.get("summary","").strip()
        set_shape_text(slide, "About_Summary",
                       summary if summary else "—", size=8)

        # ── Education ─────────────────────────────────────────────────────────
        set_shape_text(slide, "Education_Content",
                       format_education(data.get("education",[])), size=8)

        # ── Languages ─────────────────────────────────────────────────────────
        set_shape_text(slide, "Languages_Content",
                       ", ".join(data.get("languages",[])) or "—", size=8)

        # ── Hobbies ───────────────────────────────────────────────────────────
        set_shape_text(slide, "Hobbies_Content",
                       ", ".join(data.get("hobbies",[])) or "—", size=8)

        # ── Skills ────────────────────────────────────────────────────────────
        set_core_skills_shape(slide, data.get("skills", {}))

        # ── Experience: bold maroon titles + coloured bullets ─────────────────
        set_experience_shape(slide, data.get("experience", []))

        # ── Quote ─────────────────────────────────────────────────────────────
        quote = personal.get("quote","").strip()
        if quote:
            set_shape_text(slide, "Quote", quote, size=9, color=SYNPULSE_ORANGE)

        out = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
        prs.save(out.name)
        print(f"✅ PPTX saved → {out.name}")
        return out.name
    except Exception:
        print("\n❌ FAILED"); traceback.print_exc(); raise



